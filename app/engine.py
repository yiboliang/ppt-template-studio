"""
模板引擎核心：把 (theme JSON + layout JSON + 用户 deck spec) 渲染成 .pptx

设计要点
--------
1. Theme 和 Layout 完全解耦：
   - Layout 只描述"版式结构"——每个占位符在幻灯片上的相对坐标 (0~1 的比例矩形)
     以及它应该使用哪个"样式令牌" (style token，如 title / h2 / body / caption)。
   - Theme 只描述"视觉规范"——配色、字体、字号，这些值通过样式令牌被 Layout 引用。
   - 因此同一套 Layout（内容结构）可以自由套用任意 Theme（视觉风格），
     这正是"可复用模板"的核心：结构与风格正交。
2. Deck spec（用户输入）只关心"要哪个 Layout + 填什么内容"，
   不需要知道任何坐标或颜色细节。
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE_DIR = Path(__file__).parent
THEMES_DIR = BASE_DIR / "themes"
LAYOUTS_DIR = BASE_DIR / "layouts"

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


class TemplateError(ValueError):
    """Schema 校验 / 渲染阶段的错误，会被 API 层转换成 4xx 响应。"""


# --------------------------------------------------------------------------
# Schema 加载
# --------------------------------------------------------------------------

def _load_json_dir(dir_path: Path) -> dict[str, dict]:
    items: dict[str, dict] = {}
    for f in sorted(dir_path.glob("*.json")):
        data = json.loads(f.read_text(encoding="utf-8"))
        items[data["id"]] = data
    return items


def list_themes() -> list[dict]:
    themes = _load_json_dir(THEMES_DIR)
    return [{"id": t["id"], "display_name": t["display_name"]} for t in themes.values()]


def list_layouts() -> list[dict]:
    layouts = _load_json_dir(LAYOUTS_DIR)
    out = []
    for l in layouts.values():
        out.append({
            "id": l["id"],
            "name": l["name"],
            "description": l.get("description", ""),
            "placeholders": [
                {"key": p["key"], "type": p["type"], "required": p.get("required", False)}
                for p in l["placeholders"]
            ],
        })
    return out


def load_theme(theme_id: str) -> dict:
    themes = _load_json_dir(THEMES_DIR)
    if theme_id not in themes:
        raise TemplateError(f"未知主题: {theme_id}，可选: {list(themes)}")
    return themes[theme_id]


def load_layout(layout_id: str) -> dict:
    layouts = _load_json_dir(LAYOUTS_DIR)
    if layout_id not in layouts:
        raise TemplateError(f"未知版式: {layout_id}，可选: {list(layouts)}")
    return layouts[layout_id]


# --------------------------------------------------------------------------
# 样式令牌 -> 具体字体/字号/颜色
# --------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor.from_string(hex_str.upper())


def _resolve_style(theme: dict, style_token: str) -> dict:
    """把 title/subtitle/h2/body/caption/quote/index 等 token 映射成具体渲染参数。"""
    colors = theme["colors"]
    sizes = theme["fonts"]["sizes"]
    heading_font = theme["fonts"]["heading_zh"][0]
    body_font = theme["fonts"]["body_zh"][0]

    table = {
        "title": {"font": heading_font, "size": sizes["title"], "bold": True, "color": colors["text_primary"]},
        "subtitle": {"font": body_font, "size": sizes["subtitle"], "bold": False, "color": colors["text_secondary"]},
        "h2": {"font": heading_font, "size": sizes["h2"], "bold": True, "color": colors["text_primary"]},
        "body": {"font": body_font, "size": sizes["body"], "bold": False, "color": colors["text_primary"]},
        "caption": {"font": body_font, "size": sizes["caption"], "bold": False, "color": colors["text_secondary"]},
        "quote": {"font": heading_font, "size": sizes["h2"], "bold": False, "color": colors["text_primary"]},
        "index": {"font": heading_font, "size": sizes["title"], "bold": True, "color": colors["accent"]},
    }
    if style_token not in table:
        raise TemplateError(f"未知样式令牌: {style_token}")
    return table[style_token]


# --------------------------------------------------------------------------
# 几何：fractional rect -> EMU
# --------------------------------------------------------------------------

def _rect_to_emu(rect: list[float], slide_w_in: float, slide_h_in: float):
    x, y, w, h = rect
    return (
        Inches(x * slide_w_in),
        Inches(y * slide_h_in),
        Inches(w * slide_w_in),
        Inches(h * slide_h_in),
    )


# --------------------------------------------------------------------------
# 渲染
# --------------------------------------------------------------------------

def _set_background(slide, color_hex: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)


def _add_textbox(slide, rect_emu, align: str, style: dict, text: str):
    left, top, width, height = rect_emu
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(style["size"])
    run.font.bold = style["bold"]
    run.font.name = style["font"]
    run.font.color.rgb = _hex_to_rgb(style["color"])
    return box


def _add_bullet_list(slide, rect_emu, align: str, style: dict, items: list[str], numbered: bool, accent_hex: str):
    left, top, width, height = rect_emu
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)
        p.space_after = Pt(style["size"] * 0.6)
        prefix = f"{i + 1:02d}  " if numbered else "•  "
        run_prefix = p.add_run()
        run_prefix.text = prefix
        run_prefix.font.size = Pt(style["size"])
        run_prefix.font.bold = True
        run_prefix.font.name = style["font"]
        run_prefix.font.color.rgb = _hex_to_rgb(accent_hex)

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(style["size"])
        run_text.font.bold = False
        run_text.font.name = style["font"]
        run_text.font.color.rgb = _hex_to_rgb(style["color"])
    return box


def _add_cover_accent_bar(slide, theme: dict):
    slide_w_in = theme["layout"]["slide_width_in"]
    slide_h_in = theme["layout"]["slide_height_in"]
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.30 * slide_h_in), Inches(0.06 * slide_w_in), Inches(0.02 * slide_h_in)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(theme["colors"]["accent"])
    bar.line.fill.background()


def _add_page_number(slide, theme: dict, page_no: int):
    slide_w_in = theme["layout"]["slide_width_in"]
    slide_h_in = theme["layout"]["slide_height_in"]
    box = slide.shapes.add_textbox(
        Inches(slide_w_in - 1.0), Inches(slide_h_in - 0.5), Inches(0.8), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page_no)
    run.font.size = Pt(theme["fonts"]["sizes"]["caption"])
    run.font.name = theme["fonts"]["body_en"][0]
    run.font.color.rgb = _hex_to_rgb(theme["colors"]["text_secondary"])


def _render_slide(prs: Presentation, theme: dict, layout: dict, content: dict, page_no: int):
    blank_layout = prs.slide_layouts[6]  # 6 = 完全空白版式
    slide = prs.slides.add_slide(blank_layout)

    bg_key = layout.get("background", "background")
    _set_background(slide, theme["colors"][bg_key])

    slide_w_in = theme["layout"]["slide_width_in"]
    slide_h_in = theme["layout"]["slide_height_in"]

    for ph in layout["placeholders"]:
        key = ph["key"]
        value = content.get(key)
        if value in (None, "", []):
            if ph.get("required"):
                raise TemplateError(f"版式 '{layout['id']}' 缺少必填字段: {key}")
            continue

        style = _resolve_style(theme, ph["style"])
        rect_emu = _rect_to_emu(ph["rect"], slide_w_in, slide_h_in)
        align = ph.get("align", "left")

        if ph["type"] == "text":
            _add_textbox(slide, rect_emu, align, style, str(value))
        elif ph["type"] == "bullet_list":
            items = value if isinstance(value, list) else [str(value)]
            _add_bullet_list(slide, rect_emu, align, style, items, numbered=False,
                              accent_hex=theme["colors"]["accent"])
        elif ph["type"] == "numbered_list":
            items = value if isinstance(value, list) else [str(value)]
            _add_bullet_list(slide, rect_emu, align, style, items, numbered=True,
                              accent_hex=theme["colors"]["accent"])
        else:
            raise TemplateError(f"不支持的占位符类型: {ph['type']}")

    if layout["id"] == "cover" and theme.get("decor", {}).get("cover_accent_bar"):
        _add_cover_accent_bar(slide, theme)

    if theme.get("decor", {}).get("page_number") and layout["id"] not in ("cover", "closing"):
        _add_page_number(slide, theme, page_no)

    return slide


def render_deck(deck_spec: dict) -> Presentation:
    """
    deck_spec = {
        "theme": "huawei_light" | "apple_dark",
        "slides": [
            {"layout": "cover", "content": {"title": "...", "subtitle": "...", "date": "..."}},
            {"layout": "content_bullets", "content": {"title": "...", "bullets": ["...", "..."]}},
            ...
        ]
    }
    """
    if "theme" not in deck_spec:
        raise TemplateError("deck_spec 缺少 'theme' 字段")
    if "slides" not in deck_spec or not isinstance(deck_spec["slides"], list) or not deck_spec["slides"]:
        raise TemplateError("deck_spec 缺少非空的 'slides' 数组")

    theme = load_theme(deck_spec["theme"])

    prs = Presentation()
    prs.slide_width = Inches(theme["layout"]["slide_width_in"])
    prs.slide_height = Inches(theme["layout"]["slide_height_in"])

    for i, slide_spec in enumerate(deck_spec["slides"], start=1):
        if "layout" not in slide_spec:
            raise TemplateError(f"第 {i} 页缺少 'layout' 字段")
        layout = load_layout(slide_spec["layout"])
        content = slide_spec.get("content", {})
        _render_slide(prs, theme, layout, content, page_no=i)

    return prs
